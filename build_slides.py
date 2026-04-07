"""
build_slides.py
Generates POP_PREDICT_Slides.pptx — 3-minute AUGS conference presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
import copy

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
NAVY     = RGBColor(0x2E, 0x40, 0x57)
TEAL     = RGBColor(0x1A, 0x6B, 0x8A)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF5, 0xF7, 0xFA)
GRAY     = RGBColor(0x55, 0x65, 0x72)
GREEN_BG = RGBColor(0xC6, 0xEF, 0xCE)
GREEN_FG = RGBColor(0x37, 0x56, 0x23)
YELLOW_BG= RGBColor(0xFF, 0xEB, 0x9C)
YELLOW_FG= RGBColor(0x7D, 0x4F, 0x00)
RED_BG   = RGBColor(0xFF, 0xC7, 0xCE)
RED_FG   = RGBColor(0x9C, 0x00, 0x06)
LIGHT_TEAL = RGBColor(0xD6, 0xEE, 0xF5)

OUTPUT = "POP_PREDICT_Slides.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # fully blank


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=24, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_multiline_text(slide, lines, left, top, width, height,
                       font_size=20, color=WHITE, line_spacing_pt=None):
    """lines = list of (text, bold, size_override)"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, size_override) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size_override if size_override else font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
        if line_spacing_pt:
            p.line_spacing = Pt(line_spacing_pt)
    return txBox


def add_score_row(slide, left, top, width, height,
                  factor, criteria, pts, bg, fg):
    """One row of the score card."""
    col_widths = [width * 0.36, width * 0.50, width * 0.14]
    x = left
    texts = [factor, criteria, pts]
    for i, (txt, cw) in enumerate(zip(texts, col_widths)):
        r = add_rect(slide, x, top, cw - 0.01, height, fill_color=bg)
        r.line.fill.background()
        tb = add_text(slide, txt, x + 0.04, top + 0.01, cw - 0.08, height - 0.02,
                      font_size=10.5, bold=(i == 0 and factor != ""),
                      color=fg, align=PP_ALIGN.LEFT if i < 2 else PP_ALIGN.CENTER)
        x += cw


def slide_number(slide, n, total, color=GRAY):
    add_text(slide, f"{n} / {total}", 12.3, 7.15, 0.9, 0.3,
             font_size=9, color=color, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# SLIDE 1 — Title
# ---------------------------------------------------------------------------
s1 = prs.slides.add_slide(blank_layout)

# Full navy background
add_rect(s1, 0, 0, 13.333, 7.5, fill_color=NAVY)

# Teal accent bar left
add_rect(s1, 0, 0, 0.18, 7.5, fill_color=TEAL)

# Top accent line
add_rect(s1, 0.18, 2.3, 13.15, 0.06, fill_color=TEAL)

# Title
add_text(s1, "POP-PREDICT Score", 0.5, 1.0, 12.5, 1.2,
         font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_text(s1,
         "A Point-Based Clinical Risk Calculator\nfor POP Surgical Recurrence",
         0.5, 2.55, 12.5, 1.0,
         font_size=24, bold=False, color=LIGHT_TEAL, align=PP_ALIGN.CENTER)

# Risk tier color strip
tier_colors = [GREEN_BG, YELLOW_BG, RED_BG]
tier_labels = ["LOW RISK", "MODERATE RISK", "HIGH RISK"]
tier_fg     = [GREEN_FG,  YELLOW_FG,        RED_FG]
strip_w = 3.5
strip_x = (13.333 - strip_w * 3) / 2
for i, (bg, label, fg) in enumerate(zip(tier_colors, tier_labels, tier_fg)):
    add_rect(s1, strip_x + i * strip_w, 3.8, strip_w - 0.05, 0.55, fill_color=bg)
    add_text(s1, label, strip_x + i * strip_w, 3.82, strip_w - 0.05, 0.5,
             font_size=14, bold=True, color=fg, align=PP_ALIGN.CENTER)

# Authors + conference
add_text(s1, "Sullivan M  ·  Kim S  ·  El Haraki A",
         0.5, 4.65, 12.5, 0.5,
         font_size=16, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, "Atrium Health Wake Forest Baptist  ·  AUGS 2025",
         0.5, 5.1, 12.5, 0.4,
         font_size=13, italic=True, color=LIGHT_TEAL, align=PP_ALIGN.CENTER)

slide_number(s1, 1, 7, color=LIGHT_TEAL)


# ---------------------------------------------------------------------------
# SLIDE 2 — The Problem
# ---------------------------------------------------------------------------
s2 = prs.slides.add_slide(blank_layout)
add_rect(s2, 0, 0, 13.333, 7.5, fill_color=OFFWHITE)
add_rect(s2, 0, 0, 13.333, 1.1, fill_color=NAVY)
add_rect(s2, 0, 1.1, 13.333, 0.06, fill_color=TEAL)

add_text(s2, "The Problem", 0.4, 0.15, 12.5, 0.9,
         font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Big stat boxes
stats = [
    ("1 in 2", "parous women\naffected by POP"),
    ("200,000+", "prolapse surgeries\nper year in the U.S."),
    ("15–40%", "recurrence risk\ndepending on procedure"),
]
box_w = 3.6
box_gap = 0.4
total_w = box_w * 3 + box_gap * 2
start_x = (13.333 - total_w) / 2

for i, (number, label) in enumerate(stats):
    bx = start_x + i * (box_w + box_gap)
    add_rect(s2, bx, 1.5, box_w, 2.8, fill_color=NAVY)
    add_text(s2, number, bx, 1.7, box_w, 1.4,
             font_size=52, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_text(s2, label, bx, 2.9, box_w, 1.2,
             font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

# Bottom statement
add_rect(s2, 0.5, 4.7, 12.333, 0.8, fill_color=RED_BG)
add_text(s2,
         "Wide variability in recurrence risk — yet counseling remains one-size-fits-all.",
         0.6, 4.72, 12.1, 0.75,
         font_size=17, bold=True, color=RED_FG, align=PP_ALIGN.CENTER)

add_text(s2,
         "Rates vary by surgical approach, anatomy, and patient characteristics — but no tool exists to quantify individual risk preoperatively.",
         0.5, 5.65, 12.333, 0.9,
         font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

slide_number(s2, 2, 7)


# ---------------------------------------------------------------------------
# SLIDE 3 — The Gap & Our Approach
# ---------------------------------------------------------------------------
s3 = prs.slides.add_slide(blank_layout)
add_rect(s3, 0, 0, 13.333, 7.5, fill_color=OFFWHITE)
add_rect(s3, 0, 0, 13.333, 1.1, fill_color=NAVY)
add_rect(s3, 0, 1.1, 13.333, 0.06, fill_color=TEAL)

add_text(s3, "Medicine Already Does This Well", 0.4, 0.15, 12.5, 0.9,
         font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Three existing calculators as comparison boxes
calcs = [
    ("Caprini Score", "VTE risk\nafter surgery", "Validated globally\nfor surgical patients"),
    ("CHA₂DS₂-VASc", "Stroke risk\nin atrial fibrillation", "Guides anticoagulation\ndecisions worldwide"),
    ("VBAC Calculator", "Success probability\nfor trial of labor", "Used routinely in\nOB counseling"),
]
box_w = 3.5
box_gap = 0.54
start_x = (13.333 - (box_w * 3 + box_gap * 2)) / 2

for i, (name, use, note) in enumerate(calcs):
    bx = start_x + i * (box_w + box_gap)
    add_rect(s3, bx, 1.4, box_w, 2.7, fill_color=LIGHT_TEAL)
    add_rect(s3, bx, 1.4, box_w, 0.55, fill_color=TEAL)
    add_text(s3, name, bx, 1.42, box_w, 0.52,
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s3, use,  bx, 2.05, box_w, 0.9,
             font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s3, note, bx, 2.85, box_w, 0.9,
             font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

# Arrow down
add_text(s3, "↓", 5.9, 4.2, 1.5, 0.6,
         font_size=32, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Gap statement
add_rect(s3, 1.5, 4.85, 10.333, 0.75, fill_color=NAVY)
add_text(s3, "No equivalent tool exists for POP surgical recurrence",
         1.6, 4.87, 10.1, 0.7,
         font_size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s3,
         "We addressed this gap with a validated point-based score derived from individual patient data across 4 prospective trials.",
         1.0, 5.75, 11.333, 0.75,
         font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

slide_number(s3, 3, 7)


# ---------------------------------------------------------------------------
# SLIDE 4 — How We Built It (Methods)
# ---------------------------------------------------------------------------
s4 = prs.slides.add_slide(blank_layout)
add_rect(s4, 0, 0, 13.333, 7.5, fill_color=OFFWHITE)
add_rect(s4, 0, 0, 13.333, 1.1, fill_color=NAVY)
add_rect(s4, 0, 1.1, 13.333, 0.06, fill_color=TEAL)

add_text(s4, "How We Built It", 0.4, 0.15, 12.5, 0.9,
         font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Flow diagram: Data → Model → Score → Calculator
steps = [
    ("4 Trials\n752 Women", "ALTIS · BEST\nEPACT · SASS", TEAL, WHITE),
    ("Logistic\nRegression", "AIC backward\nselection", NAVY, WHITE),
    ("β → Points", "Coefficients\nscaled to integers", TEAL, WHITE),
    ("POP-PREDICT\nScore", "0–7 pts\nReady for clinic", GREEN_FG, WHITE),
]
arrow_gap = 0.25
box_w = 2.7
box_h = 2.2
total = box_w * 4 + arrow_gap * 3
sx = (13.333 - total) / 2
sy = 1.6

for i, (title, sub, bg, fg) in enumerate(steps):
    bx = sx + i * (box_w + arrow_gap)
    add_rect(s4, bx, sy, box_w, box_h, fill_color=bg)
    add_text(s4, title, bx, sy + 0.2, box_w, 0.9,
             font_size=17, bold=True, color=fg, align=PP_ALIGN.CENTER)
    add_text(s4, sub, bx, sy + 1.05, box_w, 0.9,
             font_size=12, color=fg, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(s4, "→", bx + box_w, sy + box_h/2 - 0.25, arrow_gap + 0.1, 0.5,
                 font_size=26, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Why logistic regression for a point score?
add_rect(s4, 0.5, 4.15, 12.333, 2.9, fill_color=LIGHT_TEAL)
add_text(s4, "Why Logistic Regression for a Point Score?",
         0.7, 4.2, 12.0, 0.5,
         font_size=16, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

bullets = [
    "β coefficients = log-odds contribution of each predictor — directly proportional to recurrence risk",
    "Scaling β values to integers preserves the relative weights between predictors",
    "This method is how all major clinical scores are derived: Caprini, Framingham, CHA₂DS₂-VASc",
    "Result: a score clinicians can calculate in 30 seconds without a computer",
]
for j, b in enumerate(bullets):
    add_text(s4, f"•  {b}", 0.8, 4.75 + j * 0.54, 12.0, 0.52,
             font_size=12.5, color=NAVY, align=PP_ALIGN.LEFT)

slide_number(s4, 4, 7)


# ---------------------------------------------------------------------------
# SLIDE 5 — The Score Card
# ---------------------------------------------------------------------------
s5 = prs.slides.add_slide(blank_layout)
add_rect(s5, 0, 0, 13.333, 7.5, fill_color=OFFWHITE)
add_rect(s5, 0, 0, 13.333, 1.1, fill_color=NAVY)
add_rect(s5, 0, 1.1, 13.333, 0.06, fill_color=TEAL)

add_text(s5, "The POP-PREDICT Score", 0.4, 0.15, 12.5, 0.9,
         font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Score card left column
card_left = 0.5
card_top  = 1.35
row_h     = 0.44
card_w    = 6.3

# Header row
add_rect(s5, card_left, card_top, card_w, row_h, fill_color=NAVY)
cols = [card_w * 0.36, card_w * 0.50, card_w * 0.14]
cx = card_left
for txt, cw in zip(["Risk Factor", "Criteria", "Pts"], cols):
    add_text(s5, txt, cx + 0.04, card_top + 0.06, cw - 0.08, row_h - 0.1,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.LEFT if txt != "Pts" else PP_ALIGN.CENTER)
    cx += cw

score_rows_data = [
    ("Surgical Approach", "Colpocleisis",                          "0",  LIGHT_TEAL, NAVY, True),
    ("",                  "Sacrocolpopexy",                        "1",  WHITE,      NAVY, False),
    ("",                  "Native Tissue Repair (USLS / SSLF)",    "2",  WHITE,      NAVY, False),
    ("Ba Point (Ant. Wall)","< 0 cm",                              "0",  LIGHT_TEAL, NAVY, True),
    ("",                  "0 – 3 cm",                              "1",  WHITE,      NAVY, False),
    ("",                  "≥ 4 cm",                                "2",  WHITE,      NAVY, False),
    ("Genital Hiatus",    "≤ 3 cm",                                "0",  LIGHT_TEAL, NAVY, True),
    ("",                  "4 – 5 cm",                              "1",  WHITE,      NAVY, False),
    ("",                  "≥ 6 cm",                                "2",  WHITE,      NAVY, False),
    ("Hormonal Status",   "Pre-menopausal or on estrogen",         "0",  LIGHT_TEAL, NAVY, True),
    ("",                  "Postmenopausal, no HRT",                "1",  WHITE,      NAVY, False),
]

for idx, (factor, criteria, pts, bg, fg, is_hdr) in enumerate(score_rows_data):
    ty = card_top + row_h + idx * row_h
    add_score_row(s5, card_left, ty, card_w, row_h, factor, criteria, pts, bg, fg)

# Total score bar
ty_total = card_top + row_h + len(score_rows_data) * row_h
add_rect(s5, card_left, ty_total, card_w, row_h, fill_color=NAVY)
add_text(s5, "TOTAL SCORE", card_left + 0.05, ty_total + 0.06, card_w * 0.5, row_h - 0.1,
         font_size=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s5, "_____ / 7 pts", card_left + card_w * 0.45, ty_total + 0.06, card_w * 0.5, row_h - 0.1,
         font_size=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# RIGHT: Risk tiers + key stats
right_x = card_left + card_w + 0.4
right_w = 13.333 - right_x - 0.4

# AUC callout
add_rect(s5, right_x, 1.35, right_w, 1.0, fill_color=NAVY)
add_text(s5, "AUC  0.67", right_x, 1.4, right_w, 0.55,
         font_size=36, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_text(s5, "validated across 4 trials", right_x, 1.87, right_w, 0.4,
         font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

# Tier cards
tiers = [
    ("0–2 pts", "LOW RISK",      "8–17%",  GREEN_BG,  GREEN_FG),
    ("3–4 pts", "MODERATE RISK", "21–28%", YELLOW_BG, YELLOW_FG),
    ("5–7 pts", "HIGH RISK",     "40–45%", RED_BG,    RED_FG),
]
tier_h = 1.1
tier_top = 2.6
for i, (score, label, pct, bg, fg) in enumerate(tiers):
    ty2 = tier_top + i * (tier_h + 0.12)
    add_rect(s5, right_x, ty2, right_w, tier_h, fill_color=bg)
    add_text(s5, label, right_x, ty2 + 0.04, right_w, 0.42,
             font_size=15, bold=True, color=fg, align=PP_ALIGN.CENTER)
    add_text(s5, f"{score}  →  {pct} recurrence", right_x, ty2 + 0.52, right_w, 0.45,
             font_size=13, color=fg, align=PP_ALIGN.CENTER)

slide_number(s5, 5, 7)


# ---------------------------------------------------------------------------
# SLIDE 6 — The Calculator (live tool)
# ---------------------------------------------------------------------------
s6 = prs.slides.add_slide(blank_layout)
add_rect(s6, 0, 0, 13.333, 7.5, fill_color=OFFWHITE)
add_rect(s6, 0, 0, 13.333, 1.1, fill_color=NAVY)
add_rect(s6, 0, 1.1, 13.333, 0.06, fill_color=TEAL)

add_text(s6, "Try It Now", 0.4, 0.15, 12.5, 0.9,
         font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Left: mock calculator UI
mock_x = 0.5
mock_w = 7.2
mock_top = 1.35

add_rect(s6, mock_x, mock_top, mock_w, 5.7, fill_color=WHITE)
add_rect(s6, mock_x, mock_top, mock_w, 0.55, fill_color=NAVY)
add_text(s6, "POP-PREDICT Score Calculator",
         mock_x + 0.1, mock_top + 0.06, mock_w - 0.2, 0.44,
         font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Four input fields (mock)
inputs = [
    ("Surgical Approach",   "Sacrocolpopexy — 1 pt"),
    ("Baseline Ba Point",   "0 to 3 cm — 1 pt"),
    ("Genital Hiatus",      "4 to 5 cm — 1 pt"),
    ("Hormonal Status",     "Postmenopausal, no HRT — 1 pt"),
]
for j, (label, val) in enumerate(inputs):
    iy = mock_top + 0.75 + j * 0.82
    add_text(s6, label, mock_x + 0.15, iy, mock_w * 0.4, 0.35,
             font_size=11, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    add_rect(s6, mock_x + 0.15, iy + 0.33, mock_w - 0.3, 0.38,
             fill_color=OFFWHITE, line_color=TEAL, line_width=0.8)
    add_text(s6, val, mock_x + 0.25, iy + 0.34, mock_w - 0.5, 0.35,
             font_size=11, color=GRAY, align=PP_ALIGN.LEFT)

# Score result box
add_rect(s6, mock_x + 0.15, mock_top + 4.1, mock_w - 0.3, 1.5, fill_color=NAVY)
add_text(s6, "TOTAL SCORE", mock_x + 0.15, mock_top + 4.15, mock_w - 0.3, 0.4,
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s6, "4 / 7", mock_x + 0.15, mock_top + 4.5, mock_w - 0.3, 0.65,
         font_size=38, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_rect(s6, mock_x + 1.5, mock_top + 5.1, mock_w - 3.0, 0.38, fill_color=YELLOW_BG)
add_text(s6, "MODERATE RISK  •  21–28%", mock_x + 1.5, mock_top + 5.12, mock_w - 3.0, 0.35,
         font_size=12, bold=True, color=YELLOW_FG, align=PP_ALIGN.CENTER)

# Right: QR + URL
right_x2 = mock_x + mock_w + 0.5
right_w2 = 13.333 - right_x2 - 0.4

add_text(s6, "Scan or Visit:", right_x2, 1.55, right_w2, 0.5,
         font_size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# QR placeholder box
add_rect(s6, right_x2 + 0.4, 2.1, right_w2 - 0.8, 2.8, fill_color=WHITE,
         line_color=NAVY, line_width=1.5)
add_text(s6, "[ QR Code ]", right_x2 + 0.4, 2.1, right_w2 - 0.8, 2.8,
         font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s6, "Insert QR code\nfor URL below", right_x2 + 0.4, 2.9, right_w2 - 0.8, 0.8,
         font_size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

add_rect(s6, right_x2, 5.1, right_w2, 0.55, fill_color=LIGHT_TEAL)
add_text(s6, "mesullivan01.github.io/\nPOP_Recurrence_Calculator",
         right_x2, 5.1, right_w2, 0.55,
         font_size=9.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_text(s6, "Free · No login · Works on any device",
         right_x2, 5.8, right_w2, 0.4,
         font_size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

add_text(s6, "Share with colleagues →\ncalculate in clinic, in 30 seconds",
         right_x2, 6.35, right_w2, 0.7,
         font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

slide_number(s6, 6, 7)


# ---------------------------------------------------------------------------
# SLIDE 7 — Conclusion
# ---------------------------------------------------------------------------
s7 = prs.slides.add_slide(blank_layout)
add_rect(s7, 0, 0, 13.333, 7.5, fill_color=NAVY)
add_rect(s7, 0, 0, 0.18, 7.5, fill_color=TEAL)
add_rect(s7, 0.18, 1.2, 13.15, 0.05, fill_color=TEAL)

add_text(s7, "Key Takeaways", 0.5, 0.25, 12.5, 0.9,
         font_size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

takeaways = [
    ("First validated point-based calculator for POP surgical recurrence",
     "Derived from 752 patients across 4 prospective trials"),
    ("4 simple inputs. 30 seconds. Actionable risk category.",
     "Surgery type · Ba point · Genital hiatus · Hormonal status"),
    ("Consistent discrimination across internal and external validation",
     "AUC 0.67 full cohort  ·  0.64 cross-validation  ·  0.55–0.66 LODO"),
    ("Supports shared decision-making before the OR",
     "Identify patients who may benefit from mesh repair or closer surveillance"),
]

for i, (main, sub) in enumerate(takeaways):
    ty = 1.45 + i * 1.35
    # Number circle
    add_rect(s7, 0.5, ty + 0.1, 0.55, 0.55, fill_color=TEAL)
    add_text(s7, str(i + 1), 0.5, ty + 0.1, 0.55, 0.55,
             font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s7, main, 1.25, ty + 0.07, 11.5, 0.45,
             font_size=17, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(s7, sub,  1.25, ty + 0.52, 11.5, 0.4,
             font_size=13, italic=False, color=LIGHT_TEAL, align=PP_ALIGN.LEFT)

slide_number(s7, 7, 7, color=LIGHT_TEAL)


# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
